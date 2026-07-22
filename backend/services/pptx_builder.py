"""Six-slide executive deck builder — content driven by question profile."""

from __future__ import annotations

from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

WHITE = RGBColor(255, 255, 255)
DARK_BG = RGBColor(0, 0, 0)
CARD_BG = RGBColor(10, 22, 40)
LIGHT_BLUE = RGBColor(147, 197, 253)
GRAY = RGBColor(148, 163, 184)
GREEN = RGBColor(52, 211, 153)
RED = RGBColor(248, 113, 113)
AMBER = RGBColor(251, 191, 36)
ACCENT = RGBColor(59, 130, 246)
DARK_GRAY = RGBColor(80, 80, 100)
BORDER = RGBColor(30, 58, 95)

_COLOR_MAP = {
    "red": RED,
    "green": GREEN,
    "amber": AMBER,
    "accent": ACCENT,
}


def _clip(text: str, max_len: int = 140) -> str:
    text = " ".join(str(text).split())
    return text if len(text) <= max_len else text[: max_len - 1] + "…"


def _set_para(p, text: str, size: int, color: RGBColor, *, bold: bool = False) -> None:
    p.text = _clip(text, 200 if size >= 12 else 120)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color


def _textbox(slide, left, top, width, height, text, size, color, *, bold=False, center=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    _set_para(p, text, size, color, bold=bold)
    if center:
        p.alignment = PP_ALIGN.CENTER
    return tf


def _bullets(slide, items: list[str], left=0.7, top=1.4, width=5.8, height=5.0, size=12):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items[:6]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_para(p, f"▸  {item}", size, GRAY)
        p.space_after = Pt(8)


def build_executive_deck(context: dict[str, Any] | None = None) -> bytes:
    ctx = context or {}
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def _bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BG
        bg.line.fill.background()

    def _header(slide, title: str):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.0))
        bar.fill.solid()
        bar.fill.fore_color.rgb = CARD_BG
        bar.line.fill.background()
        _textbox(slide, 0.7, 0.22, 12.0, 0.55, title, 20, WHITE, bold=True)

    def _footer(slide):
        _textbox(slide, 0.7, 6.95, 5.0, 0.3, "DeckGen AI · Confidential", 9, DARK_GRAY)

    # Slide 1 — Executive summary
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s1)
    _header(s1, "Executive Summary")
    _textbox(s1, 0.9, 1.55, 11.5, 1.4, ctx.get("s1_headline", ""), 22, LIGHT_BLUE, bold=True)
    _textbox(s1, 0.9, 3.15, 6.0, 1.2, ctx.get("s1_number", ""), 52, RED, bold=True)
    _textbox(s1, 0.9, 4.55, 11.0, 0.6, ctx.get("s1_sub", ""), 13, GRAY)
    _footer(s1)

    # Slide 2 — Cost breakdown + chart
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s2)
    _header(s2, "Cost Breakdown & Benchmarking")
    _bullets(s2, ctx.get("s2_bullets", []))
    chart_bg = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.25), Inches(5.8), Inches(5.2))
    chart_bg.fill.solid()
    chart_bg.fill.fore_color.rgb = CARD_BG
    chart_bg.line.color.rgb = BORDER
    _textbox(s2, 7.0, 1.45, 5.4, 0.4, ctx.get("s2_chart_title", "Budget vs Actual"), 11, LIGHT_BLUE, bold=True)
    bars = ctx.get("s2_chart_bars", [])[:5]
    max_val = max((v for _, v in bars), default=100)
    for i, (cat, val) in enumerate(bars):
        bh = max(Inches(0.15), (val / max_val) * Inches(2.6))
        bar = s2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(7.15 + i * 1.05),
            Inches(4.55) - bh,
            Inches(0.65),
            bh,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RED if val > 110 else GREEN
        bar.line.fill.background()
        _textbox(s2, 7.1 + i * 1.05, 4.65, 0.75, 0.3, cat[:10], 8, GRAY, center=True)
    _footer(s2)

    # Slide 3 — Waterfall
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s3)
    _header(s3, "Opportunity Waterfall — Cost Impact")
    waterfall = ctx.get("s3_waterfall", [])[:6]
    n = len(waterfall) or 1
    box_w = min(1.85, 11.5 / n - 0.15)
    x = 0.6
    for label, val, sev in waterfall:
        color = _COLOR_MAP.get(str(sev).lower(), ACCENT)
        box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.0), Inches(box_w), Inches(3.6))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        _textbox(s3, x + 0.08, 2.25, box_w - 0.16, 0.9, label, 10, GRAY, center=True)
        _textbox(s3, x + 0.08, 3.55, box_w - 0.16, 0.7, val, 15, color, bold=True, center=True)
        x += box_w + 0.18
    _textbox(s3, 0.7, 5.85, 12.0, 0.5, ctx.get("s3_note", ""), 11, GREEN)
    _footer(s3)

    # Slide 4 — Audience narratives
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s4)
    _header(s4, "Audience-Specific Narratives")
    audiences = ctx.get("s4_audiences", [])[:3]
    for i, (aud, narrative) in enumerate(audiences):
        lx = 0.7 + i * 4.1
        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lx), Inches(1.25), Inches(3.75), Inches(5.1))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER
        _textbox(s4, lx + 0.15, 1.45, 3.45, 0.45, aud, 14, LIGHT_BLUE, bold=True)
        box = s4.shapes.add_textbox(Inches(lx + 0.15), Inches(2.05), Inches(3.45), Inches(3.8))
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        _set_para(tf.paragraphs[0], narrative, 11, GRAY)
    _footer(s4)

    # Slide 5 — Risk matrix
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s5)
    _header(s5, "Risk Matrix")
    for col, h in enumerate(["Risk", "Likelihood", "Impact", "Status"]):
        _textbox(s5, 0.8 + col * 3.0, 1.35, 2.8, 0.35, h, 11, LIGHT_BLUE, bold=True)
    for row, (risk, lik, imp, sev) in enumerate(ctx.get("s5_risks", [])[:5]):
        y = 1.85 + row * 0.72
        color = _COLOR_MAP.get(str(sev).lower(), AMBER)
        for col, val in enumerate([risk, lik, imp, "●"]):
            _textbox(
                s5, 0.8 + col * 3.0, y, 2.8, 0.5,
                val, 10, color if col == 3 else GRAY,
            )
    _footer(s5)

    # Slide 6 — Next steps
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s6)
    _header(s6, "Next Steps")
    box = s6.shapes.add_textbox(Inches(0.9), Inches(1.35), Inches(11.5), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    for i, action in enumerate(ctx.get("s6_steps", [])[:6]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        _set_para(p, f"{i + 1}.  {action}", 13 if i == 0 else 12, WHITE if i == 0 else GRAY)
        p.space_after = Pt(10)
    _footer(s6)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# Legacy 4-slide builder used by FastAPI /api/generate-deck
def build_pptx(narrative: dict[str, Any], store, quarter: str) -> bytes:
    from .data import DataStore  # noqa: F401

    NAVY = RGBColor(0x0F, 0x17, 0x2A)
    SLATE = RGBColor(0x1E, 0x29, 0x3B)
    CYAN = RGBColor(0x38, 0xBD, 0xF8)
    MUTED = RGBColor(0x94, 0xA3, 0xB8)

    def _set_slide_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=WHITE):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.word_wrap = True
        p = frame.paragraphs[0]
        p.text = text
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def _add_bullets(slide, bullets, top=2.4):
        box = slide.shapes.add_textbox(Inches(0.9), Inches(top), Inches(11.5), Inches(4.5))
        frame = box.text_frame
        frame.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            p.text = bullet
            if p.runs:
                p.runs[0].font.size = Pt(18)
                p.runs[0].font.color.rgb = WHITE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for index, slide_data in enumerate(narrative.get("slides", [])):
        slide = prs.slides.add_slide(blank)
        _set_slide_bg(slide, NAVY)
        _add_textbox(slide, 0.9, 0.7 if index else 2.2, 11.5, 1.2, slide_data.get("title", ""), size=34 if index == 0 else 28, bold=True)
        subtitle = slide_data.get("subtitle", "")
        if subtitle:
            _add_textbox(slide, 0.9, 1.5 if index == 0 else 1.55, 11.5, 0.8, subtitle, size=18, color=MUTED)
        bullets = slide_data.get("bullets") or []
        if bullets:
            _add_bullets(slide, bullets, top=2.2 if index != 0 else 3.3)
        _add_textbox(slide, 0.9, 6.9, 4, 0.3, "DeckGen AI", size=10, color=MUTED)
    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
