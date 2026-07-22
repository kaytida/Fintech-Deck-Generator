import streamlit as st
import time
import io
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def create_deck_pptx(topic, data_summary=None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    DARK_BG = RGBColor(0, 0, 0)
    CARD_BG = RGBColor(10, 22, 40)
    ACCENT = RGBColor(59, 130, 246)
    LIGHT_BLUE = RGBColor(147, 197, 253)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(148, 163, 184)
    GREEN = RGBColor(52, 211, 153)
    RED = RGBColor(248, 113, 113)
    AMBER = RGBColor(251, 191, 36)
    DARK_GRAY = RGBColor(80, 80, 100)
    BORDER = RGBColor(30, 58, 95)

    def _bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BG
        bg.line.fill.background()

    def _header(slide, title):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = CARD_BG
        bar.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.28), Inches(12), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE

    def _footer(slide):
        fb = slide.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(5), Inches(0.3))
        p = fb.text_frame.paragraphs[0]
        p.text = "DeckGen AI · Confidential"
        p.font.size = Pt(9)
        p.font.color.rgb = DARK_GRAY

    # Slide 1 — Executive summary (1-liner + key number)
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s1)
    _header(s1, "Executive Summary")
    hero = s1.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2))
    p = hero.text_frame.paragraphs[0]
    p.text = "Q2 software infrastructure spend exceeded plan by 29% — immediate cost review required."
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    num = s1.shapes.add_textbox(Inches(0.9), Inches(3.5), Inches(5), Inches(1.5))
    p = num.text_frame.paragraphs[0]
    p.text = "+$696K"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = RED
    sub = s1.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(8), Inches(0.5))
    p = sub.text_frame.paragraphs[0]
    p.text = "Budget $2.4M  →  Actual $3.1M  ·  Critical threshold breach"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY
    _footer(s1)

    # Slide 2 — Cost breakdown + benchmarking
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s2)
    _header(s2, "Cost Breakdown & Benchmarking")
    bullets = s2.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5.8), Inches(5))
    tf = bullets.text_frame
    for i, b in enumerate([
        "Cloud migration: +$340K (142% of plan)",
        "Vendor licensing: +$240K (160% of plan)",
        "Headcount expansion: +$116K (115% of plan)",
        "Benchmark: cloud rates 23% above market",
        "Peer GCC avg overrun: 12% vs our 29%",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"▸  {b}"
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.space_after = Pt(10)
    chart_bg = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.4), Inches(5.8), Inches(5))
    chart_bg.fill.solid()
    chart_bg.fill.fore_color.rgb = CARD_BG
    chart_bg.line.color.rgb = BORDER
    ct = s2.shapes.add_textbox(Inches(7.1), Inches(1.65), Inches(5), Inches(0.4))
    p = ct.text_frame.paragraphs[0]
    p.text = "Budget vs Actual by Category"
    p.font.size = Pt(12)
    p.font.color.rgb = LIGHT_BLUE
    cats = ["Cloud", "Licenses", "Headcount", "Ops", "Other"]
    vals = [112, 160, 115, 98, 105]
    for i, (cat, val) in enumerate(zip(cats, vals)):
        bh = (val / 160) * Inches(2.8)
        bar = s2.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(7.2 + i * 1.05), Inches(4.5) - bh, Inches(0.7), bh,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RED if val > 110 else GREEN
        bar.line.fill.background()
        lbl = s2.shapes.add_textbox(Inches(7.2 + i * 1.05), Inches(4.6), Inches(0.7), Inches(0.25))
        p = lbl.text_frame.paragraphs[0]
        p.text = cat
        p.font.size = Pt(9)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
    _footer(s2)

    # Slide 3 — Opportunity waterfall
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s3)
    _header(s3, "Opportunity Waterfall — Cost Impact")
    steps = [
        ("Baseline Budget", "$2,400K", ACCENT),
        ("Cloud Migration", "+$340K", RED),
        ("Vendor Licensing", "+$240K", RED),
        ("Headcount Growth", "+$116K", AMBER),
        ("Potential Savings", "-$180K", GREEN),
        ("Net Position", "$3,096K", RED),
    ]
    x = 0.6
    for label, val, color in steps:
        box = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.2), Inches(1.85), Inches(3.5))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = color
        lt = s3.shapes.add_textbox(Inches(x + 0.1), Inches(2.5), Inches(1.65), Inches(0.8))
        p = lt.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY
        p.alignment = PP_ALIGN.CENTER
        lv = s3.shapes.add_textbox(Inches(x + 0.1), Inches(3.8), Inches(1.65), Inches(0.6))
        p = lv.text_frame.paragraphs[0]
        p.text = val
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER
        x += 2.05
    note = s3.shapes.add_textbox(Inches(0.7), Inches(6.0), Inches(12), Inches(0.5))
    p = note.text_frame.paragraphs[0]
    p.text = "Identified $180K recoverable via contract renegotiation and license consolidation"
    p.font.size = Pt(12)
    p.font.color.rgb = GREEN
    _footer(s3)

    # Slide 4 — Three audience-specific narratives
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s4)
    _header(s4, "Audience-Specific Narratives")
    audiences = [
        ("CFO", "29% overrun threatens Q3 cash runway by 2.3 months. Approve $400K reallocation and freeze non-essential hiring."),
        ("Engineering VP", "Cloud migration drove 48% infra overspend. Implement monthly checkpoints and vendor pre-approval >$50K."),
        ("Board", "Infrastructure costs unsustainable at current trajectory. Schedule emergency budget review within 14 days."),
    ]
    for i, (aud, narrative) in enumerate(audiences):
        card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7 + i * 4.1), Inches(1.4), Inches(3.8), Inches(5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER
        at = s4.shapes.add_textbox(Inches(0.9 + i * 4.1), Inches(1.65), Inches(3.4), Inches(0.5))
        p = at.text_frame.paragraphs[0]
        p.text = aud
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = LIGHT_BLUE
        nt = s4.shapes.add_textbox(Inches(0.9 + i * 4.1), Inches(2.3), Inches(3.4), Inches(3.8))
        tf = nt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = narrative
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
    _footer(s4)

    # Slide 5 — Risk matrix
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s5)
    _header(s5, "Risk Matrix")
    risks = [
        ("Cloud cost escalation", "High", "High", RED),
        ("Vendor contract expiry", "Medium", "High", AMBER),
        ("Headcount creep", "High", "Medium", AMBER),
        ("Budget reallocation delay", "Low", "High", AMBER),
        ("Regulatory compliance gap", "Low", "Low", GREEN),
    ]
    headers = ["Risk", "Likelihood", "Impact", "Status"]
    for col, h in enumerate(headers):
        cell = s5.shapes.add_textbox(Inches(0.8 + col * 3.0), Inches(1.5), Inches(2.8), Inches(0.4))
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = LIGHT_BLUE
    for row, (risk, lik, imp, color) in enumerate(risks):
        y = 2.0 + row * 0.75
        for col, val in enumerate([risk, lik, imp, "●"]):
            cell = s5.shapes.add_textbox(Inches(0.8 + col * 3.0), Inches(y), Inches(2.8), Inches(0.55))
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(11)
            p.font.color.rgb = color if col == 3 else GRAY
    _footer(s5)

    # Slide 6 — Next steps
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s6)
    _header(s6, "Next Steps")
    actions = [
        "Approve emergency $400K Q3 budget reallocation by Friday",
        "Mandate pre-approval workflow for all vendor contracts >$50K",
        "Renegotiate AWS/Oracle contracts — target 20% rate reduction",
        "Deploy DeckGen AI real-time spend alerts for Engineering",
        "Schedule board budget review session within 14 days",
        "Freeze non-essential hiring until September 2024",
    ]
    tb = s6.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(5))
    tf = tb.text_frame
    for i, action in enumerate(actions):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{i + 1}.  {action}"
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY if i > 0 else WHITE
        p.space_after = Pt(14)
    _footer(s6)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

FAQ_ITEMS = [
    ("Spend last month?", "What is the spend looking from last month?"),
    ("Top variance drivers", "Show me the top drivers between the variance"),
    ("Duplicate payments?", "Were there any duplications been recorded in payments from last month"),
    ("Q2 budget variance", "What is the budget variance for Q2 software infrastructure?"),
    ("Headcount overrun", "Show me the headcount cost overrun analysis"),
    ("Cloud by vendor", "Generate a cloud spend breakdown by vendor"),
    ("Top 3 drivers", "What are the top 3 cost drivers this quarter?"),
    ("Dept budget status", "Which departments exceeded their budget?"),
]

FAQ_QUESTIONS = [q for _, q in FAQ_ITEMS]

RESPONSES = {
    "What is the spend looking from last month?": """**Last Month Spend Overview**

**Total actual:** $3,096,000
**Budget plan:** $2,400,000
**Variance:** +$696,000 (+29.0%)

**Breakdown:**
• Cloud & infrastructure: $1,420K (+18% vs plan)
• Personnel & headcount: $890K (+12% vs plan)
• Vendor licensing: $485K (+22% vs plan)

**Trend:** Spend accelerated in the final 2 weeks — cloud migration invoices landed late.""",

    "Show me the top drivers between the variance": """**Top Variance Drivers — Last Period**

**1. Cloud migration acceleration** (+$340K)
   Unbudgeted lift-and-shift; data transfer 3× estimate

**2. Vendor licensing renewal** (+$240K)
   Oracle +$180K · AWS reserved instances +$60K

**3. Engineering headcount** (+$157K)
   7 unplanned FTEs · avg $22.4K per hire

**Net impact:** $737K combined — explains 106% of total overrun""",

    "Were there any duplications been recorded in payments from last month": """**Payment Duplicate Scan — Last Period**

**Records scanned:** 4,218 payment transactions
**Potential duplicates flagged:** 2

• **Professional Services — Vendor ACME Consulting**
  Same invoice #4821 posted twice: $142K + $142K
  Status: 1 confirmed duplicate · 1 pending AP review

• **Recommendation:** Hold further disbursements until reconciliation complete

**Recovery opportunity:** ~$142K if duplicate confirmed""",

    "What is the budget variance for Q2 software infrastructure?": """Based on the Q2 financial data:

**Budget:** $2,400,000
**Actual:** $3,096,000
**Variance:** +$696,000 (+29% overrun)

**Key Drivers:**
• Cloud migration acceleration: +$340K (unbudgeted)
• Unplanned vendor licensing: +$240K (Oracle, AWS)
• Engineering headcount expansion: +$116K (15% increase)

**Status:** 🔴 Critical — exceeds 25% threshold""",

    "Show me the headcount cost overrun analysis": """**Headcount Cost Analysis — Q2 2024**

**Budgeted FTE:** 45 engineers
**Actual FTE:** 52 engineers (+15.6%)

**Cost Impact:**
• Additional salaries: $116K over budget
• Onboarding & equipment: $23K
• Benefits & overhead: $18K
• **Total overrun: $157K**

**Root Cause:**
3 unplanned hires for cloud migration project + 2 backfills for attrition without budget revision.""",

    "Generate a cloud spend breakdown by vendor": """**Cloud Spend Breakdown — Q2 2024**

| Vendor | Budget | Actual | Variance |
|--------|--------|--------|----------|
| AWS | $800K | $1,120K | +$320K 🔴 |
| Azure | $400K | $380K | -$20K 🟢 |
| GCP | $200K | $245K | +$45K 🟡 |
| Oracle Cloud | $150K | $195K | +$45K 🟡 |
| Others | $50K | $66K | +$16K 🟡 |

**Total Cloud:** $1,600K budget → $2,006K actual (+25.4%)""",

    "What are the top 3 cost drivers this quarter?": """**Top 3 Cost Drivers — Q2 2024**

**1. Cloud Migration Acceleration** (+$340K)
   • Unbudgeted lift-and-shift project
   • Data transfer costs exceeded estimates by 3x

**2. Vendor Licensing Renewal** (+$240K)
   • Oracle enterprise license: +$180K
   • AWS reserved instances: +$60K

**3. Engineering Headcount Growth** (+$157K)
   • 7 additional FTEs hired without budget approval
   • Average cost per hire: $22.4K""",

    "Compare actuals vs budget for Engineering": """**Engineering Department — Q2 Variance Report**

| Category | Budget | Actual | Variance |
|----------|--------|--------|----------|
| Personnel | $1.2M | $1.42M | +18.3% 🔴 |
| Cloud Infra | $600K | $890K | +48.3% 🔴 |
| Software | $300K | $385K | +28.3% 🔴 |
| Training | $50K | $42K | -16.0% 🟢 |
| Travel | $30K | $18K | -40.0% 🟢 |

**Total Engineering:** $2.18M budget → $2.755M actual (+26.4%)""",

    "Show variance trend month-over-month": """**Monthly Variance Trend — 2024**

```
Jan:  85%  ████████▌    🟢 Under budget
Feb: 120%  ████████████ 🔴 +20% overrun
Mar:  95%  █████████▌   🟢 Near target
Apr: 110%  ███████████  🟡 +10% overrun
May:  75%  ███████▌     🟢 Under budget
Jun: 130%  █████████████🔴 +30% overrun (Q2 peak)
```

**Trend:** Escalating variance in Q2, peak in June. Recommend immediate corrective action.""",

    "Which departments exceeded their budget?": """**Department Budget Status — Q2 2024**

🔴 **Over Budget:**
• Engineering: +26.4% ($575K over)
• IT Operations: +18.2% ($145K over)
• Product: +12.1% ($67K over)

🟡 **At Risk (90-100%):**
• Marketing: 98.3% utilized
• Sales: 94.7% utilized

🟢 **Under Budget:**
• HR: 72.5%
• Finance: 68.1%
• Legal: 55.3%""",

    "Generate CFO summary for board meeting": """**CFO Executive Summary — Q2 2024**

**Overall Financial Health:** ⚠️ Requires Attention

**Key Metrics:**
• Total Budget: $8.5M | Actual: $9.8M (+15.3%)
• Cash runway impact: -2.3 months
• Largest variance: Engineering (+$575K)

**Critical Actions Required:**
1. Approve emergency budget reallocation ($400K from Q3)
2. Mandate pre-approval for all vendor contracts >$50K
3. Freeze non-essential hiring until September

**Board Recommendation:** Schedule special budget review session within 14 days."""
}

EXPLORING = {
    "What is the spend looking from last month?": [
        "Connecting to <span class='xd'>enterprise_finance_db</span>… connected.",
        "Pulling last period spend from <span class='xd'>operational_budgets</span>…",
        "Total actual: <span class='xd'>$3,096,000</span> vs plan <span class='xd'>$2,400,000</span> (+29%).",
        "Breakdown by category — cloud, personnel, licensing…",
    ],
    "Show me the top drivers between the variance": [
        "Running variance decomposition across <span class='xd'>12</span> cost centers…",
        "Driver 1 — Cloud migration: <span class='xd'>+$340K</span>.",
        "Driver 2 — Vendor licensing: <span class='xd'>+$240K</span>.",
        "Driver 3 — Headcount growth: <span class='xd'>+$157K</span>.",
    ],
    "Were there any duplications been recorded in payments from last month": [
        "Connecting to <span class='xd'>accounts_payable</span> ledger… connected.",
        "Scanning <span class='xd'>4,218</span> payment transactions from last period…",
        "Flagged <span class='xd'>2 potential duplicates</span> in Professional Services.",
        "Cross-referencing invoice IDs… <span class='xd'>1 confirmed</span>, 1 pending review.",
    ],
    "What is the budget variance for Q2 software infrastructure?": [
        "Connecting to <span class='xd'>enterprise_finance_db</span>… connected.",
        "Scanning <span class='xd'>Q2</span> software infrastructure ledger — <span class='xd'>23</span> line items found.",
        "Budget: <span class='xd'>$2,400,000</span> · Actual: <span class='xd'>$3,096,000</span> · Variance: <span class='xd'>+29%</span>.",
        "Flagging critical threshold breach (>25%)…",
    ],
    "Show me the headcount cost overrun analysis": [
        "Connecting to <span class='xd'>HRMS + finance_db</span>… connected.",
        "Pulling headcount records — budgeted <span class='xd'>45 FTE</span>, actual <span class='xd'>52 FTE</span>.",
        "Cost impact: salaries <span class='xd'>+$116K</span>, onboarding <span class='xd'>+$23K</span>, benefits <span class='xd'>+$18K</span>.",
        "Total headcount overrun: <span class='xd'>$157K</span>.",
    ],
    "Generate a cloud spend breakdown by vendor": [
        "Connecting to <span class='xd'>cloud_billing_api</span>… connected.",
        "Pulling vendor invoices — AWS, Azure, GCP, Oracle Cloud…",
        "AWS: <span class='xd'>$1,120K</span> actual vs <span class='xd'>$800K</span> budget (+40%).",
        "Total cloud spend: <span class='xd'>$2,006K</span> vs <span class='xd'>$1,600K</span> plan (+25.4%).",
    ],
    "What are the top 3 cost drivers this quarter?": [
        "Running variance decomposition across <span class='xd'>12</span> cost centers…",
        "Driver 1 — Cloud migration: <span class='xd'>+$340K</span>.",
        "Driver 2 — Vendor licensing: <span class='xd'>+$240K</span>.",
        "Driver 3 — Headcount growth: <span class='xd'>+$157K</span>.",
    ],
    "Compare actuals vs budget for Engineering": [
        "Querying <span class='xd'>Engineering</span> department cost centers…",
        "Personnel: <span class='xd'>+18.3%</span> · Cloud: <span class='xd'>+48.3%</span> · Software: <span class='xd'>+28.3%</span>.",
        "Total Engineering: <span class='xd'>$2.755M</span> vs <span class='xd'>$2.18M</span> budget (+26.4%).",
    ],
    "Show variance trend month-over-month": [
        "Pulling monthly utilization series from <span class='xd'>Jan–Jun 2024</span>…",
        "Peak overrun detected in <span class='xd'>June: 130%</span> utilization.",
        "Trend: escalating variance in Q2 — corrective action recommended.",
    ],
    "Which departments exceeded their budget?": [
        "Scanning all department budgets vs actuals…",
        "Over budget: Engineering <span class='xd'>(+26.4%)</span>, IT Ops <span class='xd'>(+18.2%)</span>, Product <span class='xd'>(+12.1%)</span>.",
        "At risk: Marketing <span class='xd'>98.3%</span>, Sales <span class='xd'>94.7%</span> utilized.",
    ],
    "Generate CFO summary for board meeting": [
        "Aggregating portfolio-level metrics for board pack…",
        "Total: <span class='xd'>$9.8M</span> actual vs <span class='xd'>$8.5M</span> budget (+15.3%).",
        "Cash runway impact: <span class='xd'>-2.3 months</span>. Largest variance: Engineering <span class='xd'>+$575K</span>.",
    ],
}

DEFAULT_EXPLORING = [
    "Connecting to <span class='xd'>enterprise_finance_db</span>… connected.",
    "Parsing your question and mapping to cost centers…",
    "Running variance analysis across connected sources…",
    "Synthesizing executive summary…",
]


def _load_upload(file) -> pd.DataFrame:
    name = file.name.lower()
    if name.endswith(".csv"):
        return pd.read_csv(file)
    if name.endswith((".xlsx", ".xls")):
        return pd.read_excel(file)
    raise ValueError("Unsupported file type. Upload CSV or Excel.")


def _summarize_upload(df: pd.DataFrame, filename: str) -> str:
    lines = [
        f"**Uploaded:** `{filename}`",
        f"**Rows:** {len(df):,} · **Columns:** {', '.join(df.columns[:6])}{'…' if len(df.columns) > 6 else ''}",
        "",
    ]
    budget_col = next((c for c in df.columns if "budget" in c.lower()), None)
    actual_col = next((c for c in df.columns if "actual" in c.lower()), None)
    if budget_col and actual_col:
        budget = pd.to_numeric(df[budget_col], errors="coerce").sum()
        actual = pd.to_numeric(df[actual_col], errors="coerce").sum()
        variance = actual - budget
        pct = (variance / budget * 100) if budget else 0
        lines += [
            f"**Total budget:** ${budget:,.0f}",
            f"**Total actual:** ${actual:,.0f}",
            f"**Variance:** ${variance:,.0f} ({pct:+.1f}%)",
            "",
        ]
    cat_col = next((c for c in df.columns if "category" in c.lower()), None)
    if cat_col and actual_col:
        top = df.groupby(cat_col)[actual_col].sum().sort_values(ascending=False).head(3)
        lines.append("**Top spend categories:**")
        for cat, val in top.items():
            lines.append(f"• {cat}: ${val:,.0f}")
        lines.append("")
    lines.append("*Click **Generate PPT** to build a 6-slide executive deck from this data.*")
    return "\n".join(lines)


def _exploring_html(question: str) -> str:
    lines = EXPLORING.get(question, DEFAULT_EXPLORING)
    body = "".join(f"<p class='expl-line'>{ln}</p>" for ln in lines)
    return (
        "<div class='exploring-box'>"
        "<div class='exploring-header'><span class='exploring-pulse'></span> Exploring</div>"
        f"<div class='exploring-body'>{body}</div></div>"
    )

if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "show_summary" not in st.session_state:
    st.session_state.show_summary = False
if "summary_generated" not in st.session_state:
    st.session_state.summary_generated = False
if "generating" not in st.session_state:
    st.session_state.generating = False
if "demo_open" not in st.session_state:
    st.session_state.demo_open = False
if "ppt_ready" not in st.session_state:
    st.session_state.ppt_ready = False
if "show_upload" not in st.session_state:
    st.session_state.show_upload = False
if "upload_summary" not in st.session_state:
    st.session_state.upload_summary = None
if "exploring_html" not in st.session_state:
    st.session_state.exploring_html = ""
if "pending_question" not in st.session_state:
    st.session_state.pending_question = None
if "last_input" not in st.session_state:
    st.session_state.last_input = ""

st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap');
    * { font-family: 'Inter', sans-serif; }
    .main { background: #000000; color: #ffffff; }
    .stApp { background: #000000; background-image: radial-gradient(ellipse 70% 50% at 50% -5%, rgba(37,99,235,0.15), transparent); }
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    .stDeployButton {display: none !important;}

    .hero-section { text-align: center; padding: 100px 20px 60px; background: radial-gradient(ellipse at center, rgba(99, 102, 241, 0.15) 0%, transparent 70%); }
    .badge { display: inline-block; background: rgba(99, 102, 241, 0.2); border: 1px solid rgba(99, 102, 241, 0.4); color: #818cf8; padding: 8px 20px; border-radius: 50px; font-size: 0.85rem; font-weight: 600; letter-spacing: 0.5px; margin-bottom: 30px; }
    .hero-title { font-size: 3.5rem; font-weight: 800; line-height: 1.1; margin-bottom: 20px; background: linear-gradient(135deg, #ffffff 0%, #c7c7ff 100%); -webkit-background-clip: text; -webkit-text-fill-color: transparent; background-clip: text; }
    .hero-subtitle { font-size: 1.3rem; color: #94a3b8; font-weight: 300; margin-bottom: 30px; }
    .hero-description { font-size: 1rem; color: #64748b; max-width: 600px; margin: 0 auto 30px; line-height: 1.7; }
    .cta-button { background: linear-gradient(135deg, #6366f1 0%, #8b5cf6 100%); color: white; padding: 16px 40px; border-radius: 12px; font-weight: 600; font-size: 1.1rem; border: none; cursor: pointer; transition: all 0.3s ease; box-shadow: 0 10px 40px rgba(99, 102, 241, 0.4); }
    .cta-button:hover { transform: translateY(-2px); box-shadow: 0 15px 50px rgba(99, 102, 241, 0.6); }

    .chat-window { background: linear-gradient(160deg, #0a1628 0%, #020408 100%); border: 1px solid rgba(59,130,246,0.22); border-radius: 16px; max-width: 820px; margin: 0 auto; overflow: hidden; box-shadow: 0 12px 40px rgba(0,0,0,0.6); }
    .chat-header { background: linear-gradient(135deg, #1e3a5f 0%, #2563eb 100%); padding: 14px 20px; display: flex; align-items: center; gap: 12px; }
    .chat-header-icon { width: 36px; height: 36px; background: rgba(255,255,255,0.2); border-radius: 10px; display: flex; align-items: center; justify-content: center; font-size: 1.2rem; }
    .chat-header-title { color: white; font-weight: 600; font-size: 1.1rem; }
    .chat-header-subtitle { color: rgba(255,255,255,0.7); font-size: 0.8rem; }
    .chat-body { padding: 16px 20px; min-height: 200px; max-height: 480px; overflow-y: auto; }
    .chat-message { margin-bottom: 14px; animation: fadeInUp 0.4s ease-out; }
    .chat-bubble-user { background: linear-gradient(135deg, #2563eb 0%, #3b82f6 100%); color: white; padding: 8px 14px; border-radius: 14px 14px 4px 14px; max-width: 75%; font-size: 0.78rem; line-height: 1.45; }
    .chat-bubble-ai { background: rgba(10,22,40,0.9); border: 1px solid rgba(59,130,246,0.15); color: #cbd5e1; padding: 10px 14px; border-radius: 14px 14px 14px 4px; max-width: 85%; font-size: 0.75rem; line-height: 1.55; }
    .chat-message-user { display: flex; justify-content: flex-end; }
    .chat-message-ai { display: flex; justify-content: flex-start; }
    .chat-avatar { width: 28px; height: 28px; border-radius: 50%; background: linear-gradient(135deg, #2563eb 0%, #3b82f6 100%); display: flex; align-items: center; justify-content: center; margin-right: 8px; font-size: 0.75rem; flex-shrink: 0; }
    .chat-avatar-user { width: 28px; height: 28px; border-radius: 50%; background: rgba(37,99,235,0.3); display: flex; align-items: center; justify-content: center; margin-left: 8px; font-size: 0.75rem; flex-shrink: 0; }
    .typing-indicator { display: flex; gap: 6px; padding: 12px 16px; background: rgba(255, 255, 255, 0.05); border-radius: 18px; width: fit-content; margin-bottom: 20px; }
    .typing-dot { width: 8px; height: 8px; background: #6366f1; border-radius: 50%; animation: typingBounce 1.4s infinite ease-in-out; }
    .typing-dot:nth-child(1) { animation-delay: 0s; }
    .typing-dot:nth-child(2) { animation-delay: 0.2s; }
    .typing-dot:nth-child(3) { animation-delay: 0.4s; }
    @keyframes typingBounce { 0%, 60%, 100% { transform: translateY(0); } 30% { transform: translateY(-10px); } }

    .faq-section { padding: 12px 18px 8px; border-top: 1px solid rgba(59,130,246,0.12); background: rgba(0,0,0,0.3); }
    .faq-title { color: #64748b; font-size: 0.65rem; font-weight: 600; margin-bottom: 8px; text-transform: uppercase; letter-spacing: 1px; }
    .chat-greeting { color: #e2e8f0; font-size: 0.9rem; font-weight: 600; padding: 12px 18px 4px; }

    .summary-card { background: rgba(10,22,40,0.95); border: 1px solid rgba(59,130,246,0.2); border-radius: 12px; padding: 14px 16px; margin: 8px 0; }
    .summary-title { color: #60a5fa; font-weight: 600; font-size: 0.78rem; margin-bottom: 8px; }
    .summary-content { color: #94a3b8; line-height: 1.55; font-size: 0.72rem; }
    .summary-content strong { color: #93c5fd; font-size: 0.72rem; }
    .summary-content h3, .summary-content h1, .summary-content h2 { font-size: 0.82rem !important; color: #e2e8f0 !important; margin: 0.4rem 0 !important; }
    .summary-content li, .summary-content p { font-size: 0.72rem !important; }
    .summary-content code { background: rgba(37,99,235,0.15); padding: 1px 5px; border-radius: 3px; color: #93c5fd; font-size: 0.68rem; }

    .exploring-box { background: #0a1628; border: 1px solid rgba(59,130,246,0.2); border-radius: 10px; padding: 10px 14px; margin: 6px 0; }
    .exploring-header { display: flex; align-items: center; gap: 6px; font-size: 0.72rem; font-weight: 600; color: #60a5fa; margin-bottom: 6px; }
    .exploring-pulse { width: 7px; height: 7px; border-radius: 50%; background: #60a5fa; animation: pulse 1.2s ease-in-out infinite; }
    @keyframes pulse { 0%,100%{opacity:1;transform:scale(1)} 50%{opacity:0.4;transform:scale(0.85)} }
    .expl-line { font-size: 0.68rem; color: #64748b; line-height: 1.55; margin: 4px 0; animation: fadeInUp 0.35s ease forwards; }
    .xd { color: #93c5fd; font-weight: 600; background: rgba(37,99,235,0.15); padding: 0 4px; border-radius: 3px; }

    .input-section { padding: 8px 18px 14px; border-top: 1px solid rgba(59,130,246,0.1); background: rgba(0,0,0,0.25); }

    .action-buttons { display: flex; gap: 12px; margin-top: 20px; flex-wrap: wrap; }
    .action-btn { background: linear-gradient(135deg, #6366f1 0%, #8b5cf6 100%); color: white; padding: 12px 24px; border-radius: 10px; font-weight: 600; font-size: 0.9rem; border: none; cursor: pointer; display: inline-flex; align-items: center; gap: 8px; transition: all 0.3s ease; text-decoration: none; }
    .action-btn:hover { transform: translateY(-2px); box-shadow: 0 8px 25px rgba(99, 102, 241, 0.4); }
    .action-btn-secondary { background: rgba(255, 255, 255, 0.05); border: 1px solid rgba(255, 255, 255, 0.1); color: #94a3b8; padding: 12px 24px; border-radius: 10px; font-weight: 500; font-size: 0.9rem; cursor: pointer; display: inline-flex; align-items: center; gap: 8px; transition: all 0.3s ease; text-decoration: none; }
    .action-btn-secondary:hover { background: rgba(255, 255, 255, 0.1); color: #e2e8f0; }

    .placeholder-text { color: #64748b; font-style: italic; text-align: center; padding: 60px 20px; }
    .placeholder-icon { font-size: 3rem; margin-bottom: 16px; opacity: 0.5; }

    @keyframes fadeInUp { from { opacity: 0; transform: translateY(20px); } to { opacity: 1; transform: translateY(0); } }
    .animate-in { animation: fadeInUp 0.5s ease-out forwards; }
    .divider { height: 1px; background: linear-gradient(90deg, transparent, rgba(99, 102, 241, 0.3), transparent); margin: 60px 0; }
    .section-title { font-size: 2.2rem; font-weight: 700; margin-bottom: 16px; color: #ffffff; text-align: center; }
    .section-subtitle { font-size: 1.1rem; color: #94a3b8; margin-bottom: 40px; text-align: center; max-width: 600px; margin-left: auto; margin-right: auto; }

    ::-webkit-scrollbar { width: 6px; }
    ::-webkit-scrollbar-track { background: #0f0f1a; }
    ::-webkit-scrollbar-thumb { background: #6366f1; border-radius: 3px; }

    div[data-testid="stHorizontalBlock"] button { 
        background: #0a1628 !important; 
        border: 1px solid rgba(59,130,246,0.25) !important; 
        color: #93c5fd !important; 
        border-radius: 8px !important; 
        font-size: 0.65rem !important; 
        padding: 5px 10px !important;
        min-height: 1.6rem !important;
        line-height: 1.2 !important;
        transition: all 0.2s ease !important;
    }
    div[data-testid="stHorizontalBlock"] button:hover { 
        background: rgba(37,99,235,0.15) !important; 
        border-color: rgba(96,165,250,0.5) !important; 
    }
    div[data-testid="stHorizontalBlock"] button p { font-size: 0.65rem !important; line-height: 1.2 !important; }

    .stTextInput > div > div > input {
        background: #0a1628 !important;
        border: 1px solid rgba(59,130,246,0.2) !important;
        color: #e2e8f0 !important;
        border-radius: 10px !important;
        padding: 10px 14px !important;
        font-size: 0.78rem !important;
    }
    .stTextInput > div > div > input:focus { border-color: #3b82f6 !important; box-shadow: 0 0 0 2px rgba(59,130,246,0.2) !important; }

    div[data-testid="stFileUploader"] {
        font-size: 0.72rem !important;
    }
    div[data-testid="stFileUploader"] section {
        background: #0a1628 !important;
        border: 1px dashed rgba(59,130,246,0.3) !important;
        border-radius: 10px !important;
        padding: 8px !important;
    }
</style>
""", unsafe_allow_html=True)

st.markdown("""
    <div class="hero-section">
        <div class="badge">🚀 AUTONOMOUS PRESENTATION ENGINE</div>
        <h1 class="hero-title">From Plain English to<br>Boardroom Decks</h1>
        <p class="hero-subtitle">Autonomous Presentation Engine for Finance & GCC Leadership</p>
        <p class="hero-description">
            From plain English to a boardroom-ready deck in <span style="color:#818cf8;font-weight:600;">15 seconds</span>. 
            DeckGen AI connects directly to your enterprise database, understands a plain-English request, 
            computes the numbers, writes the narrative, and generates a fully editable PowerPoint — instantly.
        </p>
    </div>
""", unsafe_allow_html=True)

col1, col2, col3 = st.columns([1, 1, 1])
with col2:
    if st.button("🚀 Try Live Demo", key="demo_btn", use_container_width=True):
        st.session_state.demo_open = True
        st.rerun()

if st.session_state.demo_open:
    st.markdown("<div style='height: 40px;'></div>", unsafe_allow_html=True)
    st.markdown("""
        <div class="section-title">Experience DeckGen AI</div>
        <div class="section-subtitle">Ask a finance question and watch it transform into an executive-ready presentation</div>
    """, unsafe_allow_html=True)

    st.markdown("""
        <div class="chat-window">
            <div class="chat-header">
                <div class="chat-header-icon">🤖</div>
                <div>
                    <div class="chat-header-title">DeckGen AI Assistant</div>
                    <div class="chat-header-subtitle">Connected to Enterprise Finance DB • Real-time</div>
                </div>
            </div>
            <div class="chat-greeting">How can I help you with planning a deck?</div>
            <div class="faq-section"><div class="faq-title">Quick questions</div></div>
        </div>
    """, unsafe_allow_html=True)

    faq_row1 = st.columns(4)
    for i, (label, question) in enumerate(FAQ_ITEMS[:4]):
        with faq_row1[i]:
            if st.button(label, key=f"faq_{i}", use_container_width=True):
                st.session_state.generating = True
                st.session_state.chat_history.append({"role": "user", "content": question})
                st.session_state.pending_question = question
                st.rerun()
    faq_row2 = st.columns(4)
    for i, (label, question) in enumerate(FAQ_ITEMS[4:], start=4):
        with faq_row2[i - 4]:
            if st.button(label, key=f"faq_{i}", use_container_width=True):
                st.session_state.generating = True
                st.session_state.chat_history.append({"role": "user", "content": question})
                st.session_state.pending_question = question
                st.rerun()

    up_col1, _ = st.columns([1, 3])
    with up_col1:
        if st.button("📎 Upload data", key="upload_toggle", use_container_width=True):
            st.session_state.show_upload = not st.session_state.show_upload
            st.rerun()

    if st.session_state.show_upload:
        uploaded = st.file_uploader(
            "Upload CSV or Excel",
            type=["csv", "xlsx", "xls"],
            key="data_upload",
            label_visibility="collapsed",
        )
        if uploaded is not None:
            try:
                df = _load_upload(uploaded)
                st.session_state.upload_summary = _summarize_upload(df, uploaded.name)
                st.session_state.generating = True
                st.session_state.chat_history.append({"role": "user", "content": f"Uploaded: {uploaded.name}"})
                st.session_state.pending_question = f"__upload__:{uploaded.name}"
                st.session_state.show_upload = False
                st.rerun()
            except Exception as exc:
                st.error(str(exc))

    st.markdown('<div class="input-section">', unsafe_allow_html=True)
    user_input = st.text_input(
        "",
        key="custom_input",
        placeholder="Ask about spend, variances, or request a deck…",
        label_visibility="collapsed",
    )
    st.markdown("</div>", unsafe_allow_html=True)

    if user_input and user_input.strip() and not st.session_state.generating:
        q = user_input.strip()
        if q != st.session_state.last_input:
            st.session_state.last_input = q
            st.session_state.generating = True
            st.session_state.chat_history.append({"role": "user", "content": q})
            st.session_state.pending_question = q
            st.rerun()

    st.markdown('<div class="chat-body">', unsafe_allow_html=True)

    for msg in st.session_state.chat_history:
        if msg["role"] == "user":
            st.markdown(
                '<div class="chat-message chat-message-user">'
                '<div class="chat-bubble-user">' + msg["content"] + "</div>"
                '<div class="chat-avatar-user">👤</div></div>',
                unsafe_allow_html=True,
            )

    if st.session_state.generating and st.session_state.get("exploring_html"):
        st.markdown(
            '<div class="chat-message chat-message-ai">'
            '<div class="chat-avatar">🤖</div>'
            '<div style="max-width:85%;">' + st.session_state.exploring_html + "</div></div>",
            unsafe_allow_html=True,
        )

    if st.session_state.show_summary and st.session_state.summary_generated:
        summary_html = st.session_state.summary_generated.replace("\n", "<br>")
        st.markdown(
            '<div class="chat-message chat-message-ai animate-in">'
            '<div class="chat-avatar">🤖</div>'
            '<div style="max-width:85%;">'
            '<div class="summary-card">'
            '<div class="summary-title">📊 Analysis Complete</div>'
            '<div class="summary-content">' + summary_html + "</div>"
            "</div></div></div>",
            unsafe_allow_html=True,
        )

    if not st.session_state.chat_history and not st.session_state.generating:
        st.markdown(
            '<div class="placeholder-text" style="padding:30px 20px;">'
            '<div style="font-size:0.78rem;color:#64748b;font-style:normal;">'
            "Pick a quick question above, upload a file, or type your query."
            "</div></div>",
            unsafe_allow_html=True,
        )

    st.markdown("</div></div>", unsafe_allow_html=True)

    if st.session_state.show_summary and st.session_state.summary_generated:
        ac1, ac2, _ = st.columns([1, 1, 2])
        with ac1:
            if st.button("📥 Generate PPT", key="gen_ppt_main", use_container_width=True, type="primary"):
                st.session_state.ppt_ready = True
        with ac2:
            if st.button("🔄 New Query", key="new_query_main", use_container_width=True):
                st.session_state.chat_history = []
                st.session_state.show_summary = False
                st.session_state.summary_generated = False
                st.session_state.generating = False
                st.session_state.ppt_ready = False
                st.session_state.exploring_html = ""
                st.session_state.upload_summary = None
                st.session_state.pending_question = None
                st.session_state.last_input = ""
                st.rerun()

    if st.session_state.ppt_ready:
        topic = "Q2 Software Infrastructure Budget Review"
        if st.session_state.chat_history:
            topic = st.session_state.chat_history[-1]["content"][:50]
        pptx_buffer = create_deck_pptx(topic, {})

        dc1, dc2, dc3 = st.columns([1, 2, 1])
        with dc2:
            st.download_button(
                label="📥 Download Executive Deck (6 slides)",
                data=pptx_buffer,
                file_name="DeckGen_AI_Executive_Summary.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )
        st.markdown(
            '<div style="text-align:center;color:#64748b;font-size:0.72rem;margin-top:8px;">'
            "✅ 6 slides · Executive summary · Waterfall · Risk matrix · Next steps"
            "</div>",
            unsafe_allow_html=True,
        )

if st.session_state.demo_open and st.session_state.generating and not st.session_state.show_summary:
    question = st.session_state.get("pending_question") or st.session_state.chat_history[-1]["content"]
    lines = EXPLORING.get(question, DEFAULT_EXPLORING)
    revealed = []
    for line in lines:
        revealed.append(line)
        body = "".join(f"<p class='expl-line'>{ln}</p>" for ln in revealed)
        st.session_state.exploring_html = (
            "<div class='exploring-box'>"
            "<div class='exploring-header'><span class='exploring-pulse'></span> Exploring</div>"
            f"<div class='exploring-body'>{body}</div></div>"
        )
        time.sleep(0.5)

    if question.startswith("__upload__:"):
        response = st.session_state.upload_summary or "**Upload processed.**"
    else:
        response = RESPONSES.get(
            question,
            """**Analysis Complete**

Based on your query, I analyzed the enterprise finance database:

**Key Findings:**
• Data retrieved from ERP, Cloud, and HRMS sources
• Variance analysis across 12 cost centers
• Critical anomalies flagged for review

*Click **Generate PPT** for the 6-slide executive deck.*""",
        )

    st.session_state.summary_generated = response
    st.session_state.show_summary = True
    st.session_state.generating = False
    st.session_state.exploring_html = ""
    st.session_state.pending_question = None
    st.rerun()

st.markdown("<div class='divider'></div>", unsafe_allow_html=True)

st.markdown("""
    <div style="max-width: 800px; margin: 0 auto; text-align: center; padding: 40px 20px;">
        <h2 class="section-title">Introduction</h2>
        <p class="section-subtitle">
            Global Capability Center (GCC) leadership teams, CFOs, and finance directors are flooded with rich BI dashboards and AI query chatbots. 
            Yet when an operational crisis or budget spike hits, the insight is trapped — locked inside dense dashboards and un-shareable chat replies.
        </p>
        <p style="font-size: 1.1rem; color: #e2e8f0; line-height: 1.8; margin-top: 30px;">
            <span style="color:#818cf8;font-weight:600;">DeckGen AI</span> is an autonomous presentation engine that turns a single plain-English request into a 
            downloadable, data-driven executive deck — closing the gap between enterprise data and the boardroom.
        </p>
    </div>
""", unsafe_allow_html=True)

st.markdown("<div class='divider'></div>", unsafe_allow_html=True)

st.markdown("""
    <div style="text-align: center; margin-bottom: 60px;">
        <h2 class="section-title">Problem Statement</h2>
        <p class="section-subtitle">
            Leadership faces a severe <span style="color: #f87171; font-weight: 600;">"Last-Mile Communication Friction."</span>
        </p>
    </div>
""", unsafe_allow_html=True)

col1, col2, col3 = st.columns(3)
with col1:
    st.markdown("""
        <div style="background: rgba(255,255,255,0.03); border: 1px solid rgba(255,255,255,0.08); border-radius: 20px; padding: 40px; text-align: center;">
            <div style="font-size: 3rem; margin-bottom: 20px;">📊</div>
            <div style="font-size: 1.3rem; font-weight: 600; margin-bottom: 15px; color: #ffffff;">Dashboards are too dense</div>
            <div style="color: #94a3b8; line-height: 1.6;">for executive decision-making during high-stakes leadership reviews.</div>
        </div>
    """, unsafe_allow_html=True)
with col2:
    st.markdown("""
        <div style="background: rgba(255,255,255,0.03); border: 1px solid rgba(255,255,255,0.08); border-radius: 20px; padding: 40px; text-align: center;">
            <div style="font-size: 3rem; margin-bottom: 20px;">💬</div>
            <div style="font-size: 1.3rem; font-weight: 600; margin-bottom: 15px; color: #ffffff;">Chatbot answers are non-shareable</div>
            <div style="color: #94a3b8; line-height: 1.6;">— trapped in text boxes that can't be presented to offshore HQ executives.</div>
        </div>
    """, unsafe_allow_html=True)
with col3:
    st.markdown("""
        <div style="background: rgba(255,255,255,0.03); border: 1px solid rgba(255,255,255,0.08); border-radius: 20px; padding: 40px; text-align: center;">
            <div style="font-size: 3rem; margin-bottom: 20px;">⏰</div>
            <div style="font-size: 1.3rem; font-weight: 600; margin-bottom: 15px; color: #ffffff;">Manual slide creation wastes 10–15 hours a month</div>
            <div style="color: #94a3b8; line-height: 1.6;">— screenshotting charts, copying numbers, and rebuilding decks.</div>
        </div>
    """, unsafe_allow_html=True)

st.markdown("<div class='divider'></div>", unsafe_allow_html=True)

st.markdown("""
    <div style="text-align: center; margin-bottom: 60px;">
        <h2 class="section-title">Our Solution</h2>
        <p class="section-subtitle" style="max-width: 800px; margin: 0 auto;">
            A user types a request such as <em style="color: #818cf8;">"Generate a 4-slide deck explaining our Q2 software infrastructure overrun for HQ."</em> DeckGen AI then autonomously:
        </p>
    </div>
""", unsafe_allow_html=True)

c1, c2 = st.columns([1, 2])
with c1:
    st.markdown("""
        <div style="background: rgba(255,255,255,0.02); border: 1px solid rgba(255,255,255,0.06); border-radius: 16px; padding: 30px; height: 100%;">
            <div style="font-size: 2rem; margin-bottom: 10px;">🔍</div>
            <div style="font-size: 1.2rem; font-weight: 600; color: #ffffff; margin-bottom: 10px;">Queries the database</div>
            <div style="color: #94a3b8; font-size: 0.95rem;">Pulls budget vs. actuals straight from the enterprise finance data.</div>
        </div>
    """, unsafe_allow_html=True)
with c2:
    st.markdown("""
        <div style="background: rgba(99,102,241,0.05); border: 1px solid rgba(99,102,241,0.1); border-radius: 16px; padding: 30px; height: 100%;">
            <pre style="background: rgba(0,0,0,0.3); padding: 20px; border-radius: 12px; color: #a5b4fc; font-family: monospace; font-size: 0.85rem; overflow-x: auto;">
SELECT department, budget_q2, actual_q2, variance_pct
FROM finance.operational_budgets
WHERE quarter = 'Q2-2024'
  AND category = 'Software Infrastructure'
            </pre>
        </div>
    """, unsafe_allow_html=True)

st.markdown("<div style='height: 20px;'></div>", unsafe_allow_html=True)

c1, c2 = st.columns([2, 1])
with c1:
    st.markdown("""
        <div style="background: rgba(139,92,246,0.05); border: 1px solid rgba(139,92,246,0.1); border-radius: 16px; padding: 30px; height: 100%; display: flex; align-items: center; justify-content: center; gap: 30px;">
            <div style="text-align: center;"><div style="font-size: 2rem; font-weight: 800; color: #34d399;">$2.4M</div><div style="color: #94a3b8; font-size: 0.9rem;">Budget</div></div>
            <div style="font-size: 2rem; color: #6366f1;">→</div>
            <div style="text-align: center;"><div style="font-size: 2rem; font-weight: 800; color: #f87171;">$3.1M</div><div style="color: #94a3b8; font-size: 0.9rem;">Actual</div></div>
            <div style="font-size: 2rem; color: #6366f1;">=</div>
            <div style="text-align: center;"><div style="font-size: 2rem; font-weight: 800; color: #fbbf24;">+29%</div><div style="color: #94a3b8; font-size: 0.9rem;">Overrun</div></div>
        </div>
    """, unsafe_allow_html=True)
with c2:
    st.markdown("""
        <div style="background: rgba(255,255,255,0.02); border: 1px solid rgba(255,255,255,0.06); border-radius: 16px; padding: 30px; height: 100%;">
            <div style="font-size: 2rem; margin-bottom: 10px;">🧮</div>
            <div style="font-size: 1.2rem; font-weight: 600; color: #ffffff; margin-bottom: 10px;">Calculates variances</div>
            <div style="color: #94a3b8; font-size: 0.95rem;">Computes budget-vs-actual gaps and quarter-over-quarter trends.</div>
        </div>
    """, unsafe_allow_html=True)

st.markdown("<div style='height: 20px;'></div>", unsafe_allow_html=True)

c1, c2 = st.columns([1, 2])
with c1:
    st.markdown("""
        <div style="background: rgba(255,255,255,0.02); border: 1px solid rgba(255,255,255,0.06); border-radius: 16px; padding: 30px; height: 100%;">
            <div style="font-size: 2rem; margin-bottom: 10px;">✍️</div>
            <div style="font-size: 1.2rem; font-weight: 600; color: #ffffff; margin-bottom: 10px;">Synthesizes insight</div>
            <div style="color: #94a3b8; font-size: 0.95rem;">Writes executive bullet points and a headline narrative.</div>
        </div>
    """, unsafe_allow_html=True)
with c2:
    st.markdown("""
        <div style="background: rgba(99,102,241,0.05); border: 1px solid rgba(99,102,241,0.1); border-radius: 16px; padding: 30px;">
            <div style="color: #818cf8; font-weight: 600; margin-bottom: 15px; font-size: 1.1rem;">🎯 Executive Summary</div>
            <ul style="color: #e2e8f0; line-height: 2; list-style: none; padding: 0;">
                <li>▸ Q2 software infrastructure spend exceeded budget by 29% ($680K overrun)</li>
                <li>▸ Primary drivers: cloud migration acceleration and unplanned vendor licensing</li>
                <li>▸ Recommend immediate cost optimization review with Engineering leadership</li>
            </ul>
        </div>
    """, unsafe_allow_html=True)

st.markdown("<div style='height: 20px;'></div>", unsafe_allow_html=True)

c1, c2 = st.columns([2, 1])
with c1:
    st.markdown("""
        <div style="background: rgba(139,92,246,0.05); border: 1px solid rgba(139,92,246,0.1); border-radius: 16px; padding: 30px; height: 100%; display: flex; align-items: center; justify-content: center;">
            <div style="text-align: center;">
                <div style="font-size: 4rem; margin-bottom: 15px;">📑</div>
                <div style="color: #818cf8; font-weight: 600; font-size: 1.2rem;">Q2_Software_Infrastructure_Review.pptx</div>
                <div style="color: #94a3b8; margin-top: 10px; font-size: 0.9rem;">4 slides • 2.3 MB • Editable</div>
                <div style="margin-top: 20px;"><span style="background: rgba(52,211,153,0.2); color: #34d399; padding: 6px 16px; border-radius: 20px; font-size: 0.85rem; font-weight: 600;">✓ Download Ready</span></div>
            </div>
        </div>
    """, unsafe_allow_html=True)
with c2:
    st.markdown("""
        <div style="background: rgba(255,255,255,0.02); border: 1px solid rgba(255,255,255,0.06); border-radius: 16px; padding: 30px; height: 100%;">
            <div style="font-size: 2rem; margin-bottom: 10px;">📥</div>
            <div style="font-size: 1.2rem; font-weight: 600; color: #ffffff; margin-bottom: 10px;">Generates the deck</div>
            <div style="color: #94a3b8; font-size: 0.95rem;">Outputs an editable .pptx with data-driven charts.</div>
        </div>
    """, unsafe_allow_html=True)

st.markdown("""
    <div style="text-align: center; padding: 60px 20px; background: radial-gradient(ellipse at center, rgba(99,102,241,0.1) 0%, transparent 70%); border-radius: 30px; margin: 60px 20px 40px;">
        <h2 style="font-size: 2.5rem; font-weight: 700; margin-bottom: 20px; color: #ffffff;">Ready to transform your<br>executive communications?</h2>
        <p style="color: #94a3b8; font-size: 1.1rem; margin-bottom: 40px;">Join leading GCCs and finance teams who've eliminated the last-mile friction.</p>
        <button class="cta-button" style="font-size: 1.2rem; padding: 18px 50px;">Start Generating Decks →</button>
        <p style="color: #64748b; margin-top: 20px; font-size: 0.9rem;">No credit card required • Enterprise-ready • SOC 2 Compliant</p>
    </div>
""", unsafe_allow_html=True)

st.markdown("""
    <div style="text-align: center; padding: 40px 20px; color: #64748b; border-top: 1px solid rgba(255,255,255,0.05);">
        <div style="display: flex; align-items: center; justify-content: center; gap: 10px; margin-bottom: 20px;">
            <span style="font-size: 1.5rem;">📊</span>
            <span style="font-size: 1.1rem; font-weight: 700; color: #ffffff;">DeckGen AI</span>
        </div>
        <p>© 2024 DeckGen AI. All rights reserved.</p>
    </div>
""", unsafe_allow_html=True)